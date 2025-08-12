from pptx import Presentation
from pathlib import Path
import pandas as pd
import uuid
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def convert_pptx_to_md(pptx_path):
    """
    Convert a PPTX file to Markdown and return it as a string.
    Extracts:
    - Slide titles and content
    - Tables (as Markdown)
    - Images (saved in 'images' folder relative to the PPTX location)
    """

    pptx_path = Path(pptx_path)
    output_dir = pptx_path.parent
    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    markdown_lines = []

    # Load the PPTX
    prs = Presentation(str(pptx_path))

    for slide_num, slide in enumerate(prs.slides, start=1):
        markdown_lines.append(f"# Slide {slide_num}\n")

        for shape in slide.shapes:
            # Extract text
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    markdown_lines.append(text + "\n")

            # Extract tables
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_data = []
                for row in shape.table.rows:
                    table_data.append([cell.text.strip() or " " for cell in row.cells])

                try:
                    df = pd.DataFrame(table_data)
                    md_table = df.to_markdown(index=False, headers=None)
                except Exception:
                    md_table = "\n".join(["| " + " | ".join(r) + " |" for r in table_data])

                markdown_lines.append(md_table + "\n")

            # Extract images
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                ext = image.ext
                image_name = f"{uuid.uuid4()}.{ext}"
                image_path = images_dir / image_name
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                markdown_lines.append(f"![Image](images/{image_name})\n")

        markdown_lines.append("\n---\n")  # Slide separator

    # Return Markdown content as a string
    return "\n".join(markdown_lines)


# For standalone testing
if __name__ == "__main__":
    pptx_file = "example.pptx"
    md_content = convert_pptx_to_md(pptx_file)
    print(md_content)
